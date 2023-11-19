import openpyxl # Excel 라이브러리
from openpyxl import load_workbook # Excel 라이브러리
import os
from pptx import Presetation # PPT 라이브러리
import sys
from wordcloud import WordCloud
from pptx.util import Inches
import pandas as pd

##### 1. 각각의 엑셀데이터를 하나의 엑셀을 만들어 시트로 이동

# 1-1. 비어있는 엑셀 만들기
wb = openpyxl.Workbook()
wb.create_sheet('car', 0)
wb.create_sheet('semiconductor', 0)
wb.create_sheet('metaverse', 0)
wb.create_sheet('battery', 0)
wb.save('product_info.xlsx')

# 1-2. 'practice/참고자료' 폴더의 Excel 파일 목록 담기
path = "practice/참고자료/"

file_list = os.listdir(path)
file_list_xls = []
for file in file_list:
    if ".xlsx" in file:
        file_list_xls.append(file)

# 1-3. Excel 파일 목록의 데이터프레임 출력
for file in file_list_xls:
    df = pd.read_excel(f"practice/참고자료/{file}")
    print(df)

# 1-4. 불러온 엑셀 파일의 내용을 product_info.xls 시트에 한 개씩 넣고 저장하기
path_dir = "practice/참고자료/"
files = os.listdir(path_dir)
file_list_xls = []
for file in files:
    if '.xls' in file:
        file_list_xls.append(file)

file_nm = "product_info.xlsx"
with pd.ExcelWriter(file_nm) as writer:
    for file_name in file_list_xls:
        df = pd.read_excel("practice/참고자료/" + file_name) # Excel 파일의 데이터프레임 조회
        df.to_excel(writer, sheet_name = file_name.replace('.xlsx', '')) # 'product_info.xlsx' 파일에 Excel 파일의 이름으로 시트를 만들어 저장한다.


##### 2. 'practice/참고자료' 폴더의 엑셀 파일 목록에서 제목과 내용을 합쳐서 워드클라우드 제작
prs = Presentation() # 파워포인트 객체 선언

# 2-1. file_list_xls에 Excel 파일 목록 담기
path = "practice/참고자료"
file_list = os.listdir(path)
file_list_xls = []
for file in file_list:
    if ".xlsx" in file:
        file_list_xls.append(file)

# 2-2. Excel 파일 목록의 데이터프레임 출력
for file in file_list_xls:
    df = pd.read_excel(f"practice/참고자료/{file}")
    print(df)

# 2-3. Excel 파일 목록의 제목과 내용을 합쳐서 워드클라우드로 제작
for file in file_list_xls:
    df = pd.read_excel(f"practice/참고자료/{file}")
    df["제목&내용"] = df["제목"] + df["내용"]

    str_unit_sum = ''
    for a in range(df.shape[0]):
        str_unit = df.loc[a, "제목&내용"]
        str_unit_sum = str_unit_sum + str_unit
    
    wc = WordCloud(font_path = "practice/MBJUA_ttf.ttf")
    wc.generate(str(str_unit_sum))
    product_name = file.replace(".xlsx", "")
    wc.to_file(f'practice/{product_name}_wordcloud_1.png')


##### 3. 'product_info.xlsx' 엑셀 파일의 시트별 파워포인트 제작
wb = load_workbook('product_info.xlsx')
sheet_names = wb.sheetnames

for sheet in sheet_names:
    prs = Presentation()

    # PPT 슬라이드 생성(텍스트용)
    title_slide_layout = prs.slide_layouts[0] # title_slide_layout 변수에 첫 번째 양식 추가
    slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가(첫 번쨰 양식)

    # PPT 내용 입력
    title = slide.shapes.title # 타이틀 텍스트 입력
    subtitle = slide.placeholders[1] # 텍스트 상자 줄 2번째 상자에 입력
    title.text = sheet
    subtitle.text = "뉴스기사 제목+내용 텍스트 기반 워드클라우드"

    # PPT 슬라이드 생성(사진용)
    img_path = f'practice/{sheet}_wordcloud_1.png'
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # PPT 이미지 첨부
    top = Inches(0) # 이미지의 시작 위치(y)
    left = Inches(0) # 이미지의 시작 위치(x)
    height = Inches(10) # 이미지의 세로 너비
    width = Inches(10 / 1.33) # 이미지의 가로 너비
    pic = slide.shapes.add_picture(img_path, left, top, height, width)
    prs.save(f'practice/{sheet}.pptx')
