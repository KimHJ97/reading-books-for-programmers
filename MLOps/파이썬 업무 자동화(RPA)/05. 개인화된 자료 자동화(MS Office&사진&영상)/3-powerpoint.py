# 3. 'product_info.xlsx' 엑셀 파일의 시트별 파워포인트 제작

# !pip install python-pptx == 0.6.21

from openpyxl import load_workbook # Excel 라이브러리
from pptx import Presentation # PopwerPoint 라이브러리
from pptx.util import Inches

wb = load_workbook('product_info.xlsx')
sheet_names = wb.sheetnames

for sheet in sheet_names:
    prs = Presentation()

    # PPT 슬라이드 생성(텍스트용)
    title_slide_layout = prs.slide_layouts[0] # title_slide_layout 변수에 첫 번째 양식 추가
    slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가(첫 번쨰 양식)

    # PPT 내용 입력
    title = slide.shapes.title # 타이틀 텍스트 입력
    subtitle = slide.placeholders[1] # 텍스트 상자 줄 2번째 상자에 입력
    title.text = sheet
    subtitle.text = "뉴스기사 제목+내용 텍스트 기반 워드클라우드"

    # PPT 슬라이드 생성(사진용)
    img_path = f'practice/{sheet}_wordcloud_1.png'
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # PPT 이미지 첨부
    top = Inches(0) # 이미지의 시작 위치(y)
    left = Inches(0) # 이미지의 시작 위치(x)
    height = Inches(10) # 이미지의 세로 너비
    width = Inches(10 / 1.33) # 이미지의 가로 너비
    pic = slide.shapes.add_picture(img_path, left, top, height, width)
    prs.save(f'practice/{sheet}.pptx')
